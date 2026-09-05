"""
Turn an already-accepted (or already-stored) resource into .docx / .pptx.

Reads Mongo only. Never generates, never judges. Citations/page hints from params and
the source preview go in the footer so the export stays traceable to the PDF.
"""

from __future__ import annotations

from io import BytesIO
import re
from typing import Dict, Tuple

from learnmate.storage import content_store, pdf_store

from . import ownership as access


def _filename_stem(resource: Dict) -> str:
    document = pdf_store.get_document(resource.get("doc_id")) if resource.get("doc_id") else None
    name = (document or {}).get("filename") or "document"
    stem = re.sub(r"\.pdf$", "", name, flags=re.IGNORECASE)
    stem = re.sub(r"[^\w\-]+", "_", stem).strip("_") or "resource"
    task = resource.get("task") or "resource"
    return f"{stem}_{task}"


def _source_line(resource: Dict) -> str:
    document = pdf_store.get_document(resource.get("doc_id")) if resource.get("doc_id") else None
    name = (document or {}).get("filename") or "source document"
    params = resource.get("params") or {}
    pages = params.get("pages")
    if isinstance(pages, list) and pages:
        return f"Source: {name} · pages {', '.join(str(p) for p in pages)}"
    if params.get("whole_document"):
        return f"Source: {name} · whole document"
    return f"Source: {name}"


def _add_markdown(paragraph, text: str) -> None:
    parts = re.split(r"(\*\*[^*]+\*\*)", str(text or ""))
    for part in parts:
        if part.startswith("**") and part.endswith("**") and len(part) >= 4:
            run = paragraph.add_run(part[2:-2])
            run.bold = True
        else:
            paragraph.add_run(part)


def _docx_bytes(resource: Dict) -> bytes:
    try:
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError as exc:
        raise ValueError(
            "Word export needs python-docx. Install it with: pip install python-docx"
        ) from exc

    doc = Document()
    title = _filename_stem(resource).replace("_", " ")
    heading = doc.add_heading(title, level=1)
    heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
    meta = doc.add_paragraph(_source_line(resource))
    meta.italic = True

    task = resource.get("task")
    content = resource.get("content")
    params = resource.get("params") or {}

    if task == "summary":
        style = params.get("summary_style") or "narrative"
        text = content if isinstance(content, str) else str(content or "")
        if style == "structured":
            for line in text.splitlines():
                line = line.strip()
                if not line:
                    continue
                line = re.sub(r"^[-*]\s+", "", line)
                paragraph = doc.add_paragraph(style="List Bullet")
                _add_markdown(paragraph, line)
        else:
            for block in re.split(r"\n\s*\n", text):
                if block.strip():
                    doc.add_paragraph(block.strip())
    elif task == "keypoints":
        for i, point in enumerate(content or [], start=1):
            paragraph = doc.add_paragraph()
            paragraph.add_run(f"{i}. ").bold = True
            _add_markdown(paragraph, str(point))
    elif task == "practice_qsn":
        for i, item in enumerate(content or [], start=1):
            q = doc.add_paragraph()
            q.add_run(f"Q{i}. ").bold = True
            q.add_run(str((item or {}).get("question") or ""))
            a = doc.add_paragraph()
            a.add_run("Answer: ").italic = True
            a.add_run(str((item or {}).get("answer") or ""))
    elif task == "mcq":
        for i, item in enumerate(content or [], start=1):
            q = doc.add_paragraph()
            q.add_run(f"Q{i}. ").bold = True
            q.add_run(str((item or {}).get("question") or ""))
            for label, option in zip("ABCD", (item or {}).get("options") or []):
                doc.add_paragraph(f"    {label}) {option}")
            ans = doc.add_paragraph()
            ans.add_run("Correct: ").italic = True
            ans.add_run(str((item or {}).get("correct_answer") or ""))
    else:
        doc.add_paragraph(str(content))

    section = doc.sections[0]
    footer = section.footer
    footer_p = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
    footer_p.text = _source_line(resource) + " · LearnMate export (verbatim; not regenerated)"

    buffer = BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


def _pptx_bytes(resource: Dict) -> bytes:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise ValueError(
            "PowerPoint export needs python-pptx. Install it with: pip install python-pptx"
        ) from exc

    prs = Presentation()
    blank = prs.slide_layouts[6]
    title_layout = prs.slide_layouts[0]
    task = resource.get("task")
    content = resource.get("content")
    source = _source_line(resource)

    def _textbox(slide, left, top, width, height, text, size=18, bold=False):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        return tf

    cover = prs.slides.add_slide(title_layout)
    cover.shapes.title.text = _filename_stem(resource).replace("_", " ")
    if cover.placeholders and len(cover.placeholders) > 1:
        cover.placeholders[1].text = source

    if task == "mcq":
        items = list(content or [])
        for i, item in enumerate(items, start=1):
            slide = prs.slides.add_slide(blank)
            _textbox(slide, 0.5, 0.3, 9, 1.4, f"{i}. {(item or {}).get('question') or ''}", 22, True)
            top = 1.9
            for label, option in zip("ABCD", (item or {}).get("options") or []):
                _textbox(slide, 0.7, top, 8.5, 0.7, f"{label}) {option}", 18)
                top += 0.85
            _textbox(slide, 0.5, 6.9, 9, 0.4, source, 10)
        key = prs.slides.add_slide(blank)
        _textbox(key, 0.5, 0.4, 9, 0.6, "Answer key", 28, True)
        lines = []
        for i, item in enumerate(items, start=1):
            lines.append(f"{i}. {(item or {}).get('correct_answer') or ''}")
        _textbox(key, 0.5, 1.2, 9, 5.2, "\n".join(lines), 16)
        _textbox(key, 0.5, 6.9, 9, 0.4, source, 10)
    elif task == "keypoints":
        points = list(content or [])
        group = 4
        for start in range(0, len(points), group):
            slide = prs.slides.add_slide(blank)
            chunk = points[start:start + group]
            _textbox(slide, 0.5, 0.3, 9, 0.5, "Key points", 24, True)
            top = 1.0
            for offset, point in enumerate(chunk):
                _textbox(slide, 0.6, top, 8.8, 1.2, f"{start + offset + 1}. {point}", 18)
                top += 1.3
            _textbox(slide, 0.5, 6.9, 9, 0.4, source, 10)
    elif task == "summary":
        slide = prs.slides.add_slide(blank)
        text = content if isinstance(content, str) else str(content or "")
        _textbox(slide, 0.5, 0.3, 9, 0.5, "Summary", 24, True)
        _textbox(slide, 0.5, 1.0, 9, 5.5, text[:1800], 16)
        _textbox(slide, 0.5, 6.9, 9, 0.4, source, 10)
    elif task == "practice_qsn":
        for i, item in enumerate(content or [], start=1):
            slide = prs.slides.add_slide(blank)
            _textbox(slide, 0.5, 0.4, 9, 1.5, f"Q{i}. {(item or {}).get('question') or ''}", 22, True)
            _textbox(slide, 0.5, 2.2, 9, 4, f"Answer: {(item or {}).get('answer') or ''}", 18)
            _textbox(slide, 0.5, 6.9, 9, 0.4, source, 10)
    else:
        slide = prs.slides.add_slide(blank)
        _textbox(slide, 0.5, 1, 9, 4, str(content)[:1500], 16)

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def export_resource(user_id: str, resource_id: str, fmt: str) -> Tuple[bytes, str, str]:
    """
    Format a stored resource. Returns (bytes, media_type, download_name).
    """
    resource = access.require_resource(user_id, resource_id)
    fmt = (fmt or "docx").strip().lower()
    if fmt not in ("docx", "pptx"):
        raise ValueError("Export format must be docx or pptx.")

    stem = _filename_stem(resource)
    if fmt == "docx":
        return (
            _docx_bytes(resource),
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            f"{stem}.docx",
        )
    return (
        _pptx_bytes(resource),
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        f"{stem}.pptx",
    )
