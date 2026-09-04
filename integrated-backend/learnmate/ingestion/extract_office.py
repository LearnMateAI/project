"""
Turn Word, PowerPoint, and LaTeX bytes into the same page list a PDF extract produces.

    [{page_index, page_number, text}, ...]

Resource generation and chat retrieve only read this list (then the cleaned/chunked
form). They do not care whether a 'page' was a PDF page, a slide, a Word slice, or a
LaTeX section.
"""

from __future__ import annotations

import re
from io import BytesIO
from typing import Dict, List

# Roughly one printed page of notes. Word has no reliable page map without Word itself.
_DOCX_PAGE_CHARS = 1800


def extract_docx(data: bytes) -> List[Dict]:
    try:
        from docx import Document
    except ImportError as exc:
        raise ValueError(
            "Word upload needs python-docx. Install it with: pip install python-docx"
        ) from exc

    try:
        document = Document(BytesIO(data))
    except Exception:
        raise ValueError(
            "This Word file could not be read. Save it as .docx and try again."
        )

    blocks: List[str] = []
    for paragraph in document.paragraphs:
        text = (paragraph.text or "").strip()
        if not text:
            continue
        if getattr(paragraph.paragraph_format, "page_break_before", False) and blocks:
            blocks.append("\f")
        blocks.append(text)

    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
            if cells:
                blocks.append(" | ".join(cells))

    if not blocks:
        return []

    pages: List[List[str]] = [[]]
    running = 0
    for block in blocks:
        if block == "\f":
            if pages[-1]:
                pages.append([])
                running = 0
            continue
        if pages[-1] and running + len(block) > _DOCX_PAGE_CHARS:
            pages.append([])
            running = 0
        pages[-1].append(block)
        running += len(block) + 1

    return _as_pages(["\n".join(part) for part in pages if part])


def extract_pptx(data: bytes) -> List[Dict]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ValueError(
            "PowerPoint upload needs python-pptx. Install it with: pip install python-pptx"
        ) from exc

    try:
        deck = Presentation(BytesIO(data))
    except Exception:
        raise ValueError(
            "This PowerPoint file could not be read. Save it as .pptx and try again."
        )

    texts = []
    for slide in deck.slides:
        parts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = (shape.text_frame.text or "").strip()
                if text:
                    parts.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame
            note_text = (notes.text or "").strip() if notes is not None else ""
            if note_text:
                parts.append(note_text)
        texts.append("\n".join(parts))
    return _as_pages(texts)


_TEX_COMMENT = re.compile(r"(?<!\\)%[^\n]*")
_TEX_SECTION = re.compile(
    r"\\(?:chapter|section|subsection|subsubsection)\*?\{([^{}]*)\}"
)
_TEX_TEXTCMD = re.compile(
    r"\\(?:textbf|textit|emph|underline|texttt|textsc|caption|title|author)\*?\{([^{}]*)\}"
)
_TEX_BEGINEND = re.compile(r"\\(?:begin|end)\{[^}]+\}")
_TEX_COMMAND = re.compile(r"\\[a-zA-Z]+\*?(?:\[[^\]]*\])?(?:\{[^{}]*\})?")


def extract_tex(data: bytes) -> List[Dict]:
    try:
        source = data.decode("utf-8")
    except UnicodeDecodeError:
        source = data.decode("latin-1")

    source = source.replace("\r\n", "\n")
    if r"\begin{document}" in source:
        source = source.split(r"\begin{document}", 1)[1]
    source = source.replace(r"\end{document}", "")
    source = _TEX_COMMENT.sub("", source)
    source = _TEX_TEXTCMD.sub(r"\1", source)

    parts = _TEX_SECTION.split(source)
    # split keeps headings in odd slots when the file uses \section{...}.
    sections: List[str] = []
    if len(parts) == 1:
        body = _strip_tex_commands(parts[0])
        if body:
            sections.append(body)
    else:
        preamble = _strip_tex_commands(parts[0])
        if preamble:
            sections.append(preamble)
        for index in range(1, len(parts), 2):
            title = (parts[index] or "").strip()
            body = _strip_tex_commands(parts[index + 1] if index + 1 < len(parts) else "")
            block = f"{title}\n\n{body}".strip() if title else body
            if block:
                sections.append(block)

    return _as_pages(sections)


def _strip_tex_commands(text: str) -> str:
    text = _TEX_BEGINEND.sub(" ", text)
    text = _TEX_COMMAND.sub(" ", text)
    text = text.replace("~", " ")
    text = re.sub(r"[{}]", " ", text)
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _as_pages(texts: List[str]) -> List[Dict]:
    pages = []
    number = 0
    for raw in texts:
        text = (raw or "").strip()
        if not text:
            continue
        pages.append({
            "page_index": number,
            "page_number": number + 1,
            "text": text,
        })
        number += 1
    return pages
