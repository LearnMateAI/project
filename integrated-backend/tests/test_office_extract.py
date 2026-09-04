"""Extractors for Word, PowerPoint, and LaTeX — no Mongo or Qdrant."""

from __future__ import annotations

import unittest
from io import BytesIO

from learnmate.ingestion.extract_office import extract_docx, extract_pptx, extract_tex
from learnmate.ingestion.formats import detect_kind
from learnmate.ingestion.validate import validate_upload


def _docx_bytes(paragraphs):
    from docx import Document

    document = Document()
    for text in paragraphs:
        document.add_paragraph(text)
    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _pptx_bytes(slides):
    from pptx import Presentation
    from pptx.util import Inches

    deck = Presentation()
    blank = deck.slide_layouts[6]
    for lines in slides:
        slide = deck.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(4))
        box.text_frame.text = "\n".join(lines)
    buffer = BytesIO()
    deck.save(buffer)
    return buffer.getvalue()


class FormatTests(unittest.TestCase):
    def test_detect_kind(self):
        self.assertEqual(detect_kind("notes.PDF"), "pdf")
        self.assertEqual(detect_kind("brief.docx"), "docx")
        self.assertEqual(detect_kind("lecture.PPTX"), "pptx")
        self.assertEqual(detect_kind("statute.tex"), "tex")

    def test_legacy_office_is_refused(self):
        with self.assertRaises(ValueError) as caught:
            detect_kind("old.ppt")
        self.assertIn(".pptx", str(caught.exception))
        with self.assertRaises(ValueError):
            detect_kind("old.doc")


class ExtractTests(unittest.TestCase):
    def test_docx_keeps_paragraph_text(self):
        data = _docx_bytes([
            "Consideration is the price of the promise.",
            "A gift promise is not a bargain.",
        ])
        pages = extract_docx(data)
        blob = " ".join(page["text"] for page in pages)
        self.assertIn("Consideration", blob)
        self.assertIn("gift promise", blob)
        kind, count = validate_upload(data, "notes.docx")
        self.assertEqual(kind, "docx")
        self.assertGreaterEqual(count, 1)

    def test_pptx_one_page_per_slide(self):
        data = _pptx_bytes([
            ["Audi alteram partem", "Hear the other side"],
            ["Ultra vires", "Outside the empowering statute"],
        ])
        pages = extract_pptx(data)
        self.assertEqual(len(pages), 2)
        self.assertIn("Audi alteram partem", pages[0]["text"])
        self.assertIn("Ultra vires", pages[1]["text"])
        kind, count = validate_upload(data, "lecture.pptx")
        self.assertEqual((kind, count), ("pptx", 2))

    def test_tex_splits_on_sections(self):
        source = r"""
\documentclass{article}
\begin{document}
\section{Consideration}
Consideration is the price of the promise.
\section{Natural justice}
Audi alteram partem requires a hearing.
\end{document}
"""
        pages = extract_tex(source.encode("utf-8"))
        self.assertGreaterEqual(len(pages), 2)
        joined = " ".join(page["text"] for page in pages)
        self.assertIn("price of the promise", joined)
        self.assertIn("Audi alteram partem", joined)
        self.assertNotIn("documentclass", joined)
        kind, count = validate_upload(source.encode("utf-8"), "notes.tex")
        self.assertEqual(kind, "tex")
        self.assertGreaterEqual(count, 2)

    def test_empty_docx_is_rejected(self):
        data = _docx_bytes([])
        with self.assertRaises(ValueError) as caught:
            validate_upload(data, "empty.docx")
        self.assertIn("No extractable text", str(caught.exception))


if __name__ == "__main__":
    unittest.main()
